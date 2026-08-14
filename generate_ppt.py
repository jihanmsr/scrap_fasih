import collections
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()

# SLIDE 1: Kesenjangan Target
slide_layout_1 = prs.slide_layouts[0] # Title slide
slide_1 = prs.slides.add_slide(slide_layout_1)
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]

title_1.text = "Analisis Kesenjangan Target UTP (ST2023 vs SE2026)"
title_1.text_frame.paragraphs[0].font.size = Pt(36)

subtitle_1.text = (
    "Target UTP Pusat (ST2023) = 843.232 Target\n"
    "Target Prelist 'Keluarga' SE2026 = 640.715 Target\n\n"
    "MENGAPA SUSUT ~200 RIBU?\n"
    "Data UTP ST2023 telah di-matching-kan oleh Pusat sebelum dilimpahkan ke SE2026.\n"
    "Sebagian besar target UTP yang sudah memiliki NIB/Izin UMKM "
    "telah dialihkan ke kategori 'OSS Perorangan' dan 'UMKM', atau dikosongkan (Blank).\n"
    "Maka, target murni UTP (rumah tangga biasa) yang dicacah adalah 640.715."
)
for p in subtitle_1.text_frame.paragraphs:
    p.font.size = Pt(18)
    p.alignment = PP_ALIGN.LEFT

# SLIDE 2: Nasib Prelist di Lapangan & Bottleneck
slide_layout_2 = prs.slide_layouts[1] # Title and Content
slide_2 = prs.slides.add_slide(slide_layout_2)
title_2 = slide_2.shapes.title
body_2 = slide_2.shapes.placeholders[1]

title_2.text = "Realisasi Lapangan & Hambatan Pengawas"

tf = body_2.text_frame
tf.text = "1. Nasib 640.715 Target Keluarga di Lapangan:"
p1 = tf.add_paragraph()
p1.text = "Berhasil Ditemukan: 434.367 (68%)"
p1.level = 1
p2 = tf.add_paragraph()
p2.text = "Gagal (Pindah/Meninggal/Tutup): 206.348 (32%)"
p2.level = 1
p3 = tf.add_paragraph()
p3.text = "Artinya: Realisasi kunjungan pencacah hampir 100%! Target gugur murni karena faktor demografi."
p3.level = 1

p4 = tf.add_paragraph()
p4.text = "\n2. Status Dokumen FASIH (Tragedi Bottleneck):"
p4.level = 0
p5 = tf.add_paragraph()
p5.text = "APPROVED BY Pengawas: 436.339 dokumen"
p5.level = 1
p6 = tf.add_paragraph()
p6.text = "SUBMITTED BY Pencacah: 157.120 dokumen (Menunggu Approval)"
p6.level = 1
p6.font.color.rgb = RGBColor(255, 0, 0)
p7 = tf.add_paragraph()
p7.text = "REJECTED BY Pengawas: 26.682 dokumen"
p7.level = 1
p8 = tf.add_paragraph()
p8.text = "DRAFT / Belum Selesai: 1.161 dokumen"
p8.level = 1

prs.save('/Users/jihanmaisaroh/scrap_fasih/Presentasi_Eksplorasi_SE2026.pptx')
print("PPT generated!")
